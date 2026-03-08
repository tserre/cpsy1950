"""Generate W8 Neural Alignment lecture deck with text + figure slides."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from PIL import Image

TEMPLATE = '/Users/tserre/Projects/teaching/cpsy1950/cpsy1950_template.pptx'
OUTPUT = '/Users/tserre/Projects/teaching/cpsy1950/presentations/W8_Neural_Alignment_Figures.pptx'
FIG_DIR = '/Users/tserre/Projects/teaching/cpsy1950/figures/week8'

LAYOUT_TITLE = 1
LAYOUT_SECTION = 2
LAYOUT_CONTENT = 3

# Content area dimensions (where figures/text body go)
AREA_LEFT = 0.34; AREA_TOP = 1.30; AREA_WIDTH = 9.32; AREA_HEIGHT = 3.60

INDIGO = RGBColor(0x4F, 0x46, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x33, 0x41, 0x55)
LIGHT_GRAY = RGBColor(0x64, 0x74, 0x8B)

NOTES_BODY_XML = '''<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/>
  <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
  <p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
  <p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>
  <a:p><a:r><a:t>{caption}</a:t></a:r></a:p></p:txBody></p:sp>'''


def remove_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def add_notes(slide, text):
    if not text:
        return
    ns = slide.notes_slide
    safe = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
    sp = etree.fromstring(NOTES_BODY_XML.format(caption=safe).encode('utf-8'))
    spTree = ns.element.find(
        './/{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    spTree.append(sp)


def add_title_slide(prs, lecture_title, subtitle):
    layout = prs.slide_layouts[LAYOUT_TITLE]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 100:
            ph.left = Inches(0.00); ph.top = Emu(int(0.04 * 914400))
            ph.width = Inches(10.00); ph.height = Emu(int(4.22 * 914400))
            tf = ph.text_frame; tf.clear()
            # Line 1: CPSY 1950
            p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run(); r1.text = "CPSY 1950"
            r1.font.name = "Outfit"; r1.font.size = Pt(20); r1.font.bold = True
            r1.font.color.rgb = INDIGO
            # Line 2: Course subtitle
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = "Deep Learning in Brains, Minds & Machines"
            r2.font.name = "Outfit"; r2.font.size = Pt(28); r2.font.bold = True
            r2.font.color.rgb = INDIGO
            # Empty line
            p_empty = tf.add_paragraph(); p_empty.alignment = PP_ALIGN.CENTER
            # Line 3: Lecture title
            p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run(); r3.text = lecture_title
            r3.font.name = "Outfit"; r3.font.size = Pt(44); r3.font.bold = True
            r3.font.color.rgb = WHITE
        elif idx == 101:
            ph.left = Inches(1.00); ph.top = Emu(int(4.35 * 914400))
            ph.width = Inches(8.00); ph.height = Emu(int(0.90 * 914400))
            tf = ph.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = subtitle
            r.font.name = "Outfit"; r.font.size = Pt(28)
            r.font.bold = False; r.font.color.rgb = GRAY


def add_section_slide(prs, title, subtitle):
    layout = prs.slide_layouts[LAYOUT_SECTION]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 101: ph.text = title
        elif idx == 102: ph.text = subtitle
    return slide


def _set_content_title(ph, title):
    """Set title placeholder on CPSY_CONTENT slide."""
    ph.left = Inches(0.00); ph.top = Emu(int(0.154 * 914400))
    ph.width = Inches(10.00); ph.height = Emu(int(0.747 * 914400))
    ph.text = title
    for para in ph.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.name = "Outfit"; run.font.size = Pt(32)
            run.font.bold = True; run.font.color.rgb = INDIGO


def _set_citation(ph, authors, journal, year):
    """Set citation placeholder — LEFT aligned."""
    ph.left = Inches(0.00); ph.top = Emu(int(5.285 * 914400))
    ph.width = Inches(10.00); ph.height = Emu(int(0.335 * 914400))
    tf = ph.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r_auth = p.add_run(); r_auth.text = authors + " "
    r_auth.font.name = "Zilla Slab Light"; r_auth.font.size = Pt(12)
    r_auth.font.bold = False; r_auth.font.italic = False; r_auth.font.color.rgb = GRAY
    r_jy = p.add_run(); r_jy.text = f"{journal} {year}"
    r_jy.font.name = "Zilla Slab Light"; r_jy.font.size = Pt(12)
    r_jy.font.bold = True; r_jy.font.italic = True; r_jy.font.color.rgb = GRAY


def fit_image_in_area(img_w, img_h, area_w, area_h):
    img_ratio = img_w / img_h
    area_ratio = area_w / area_h
    if img_ratio > area_ratio:
        w = area_w; h = w / img_ratio
    else:
        h = area_h; w = h * img_ratio
    return (area_w - w) / 2, (area_h - h) / 2, w, h


def add_figure_slide(prs, image_path, title, authors, journal, year, caption=None):
    layout = prs.slide_layouts[LAYOUT_CONTENT]
    slide = prs.slides.add_slide(layout)
    ph_pic = None
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 100: _set_content_title(ph, title)
        elif idx == 101: ph_pic = ph
        elif idx == 102: _set_citation(ph, authors, journal, year)
    
    # Remove picture placeholder, add as regular shape
    if ph_pic:
        ph_pic.element.getparent().remove(ph_pic.element)
    if os.path.exists(image_path):
        with Image.open(image_path) as im:
            img_w, img_h = im.size
        ox, oy, w, h = fit_image_in_area(img_w, img_h, AREA_WIDTH, AREA_HEIGHT)
        slide.shapes.add_picture(
            image_path,
            Emu(int((AREA_LEFT + ox) * 914400)), Emu(int((AREA_TOP + oy) * 914400)),
            Emu(int(w * 914400)), Emu(int(h * 914400)))
    add_notes(slide, caption)
    return slide


def add_text_slide(prs, title, bullets):
    """Add a text content slide.
    
    bullets: list of tuples (level, text, bold_prefix)
      level: 0 = main bullet, 1 = sub-bullet
      text: the bullet text
      bold_prefix: optional bold text before the rest (e.g., "fMRI: ")
    """
    layout = prs.slide_layouts[LAYOUT_CONTENT]
    slide = prs.slides.add_slide(layout)
    
    # Set title, remove picture placeholder and citation
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 100:
            _set_content_title(ph, title)
        elif idx == 101:
            ph.element.getparent().remove(ph.element)
        elif idx == 102:
            ph.element.getparent().remove(ph.element)
    
    # Add text body shape
    from pptx.util import Inches as In
    txBox = slide.shapes.add_textbox(
        Emu(int(0.60 * 914400)), Emu(int(1.30 * 914400)),
        Emu(int(8.80 * 914400)), Emu(int(3.80 * 914400)))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, (level, text, bold_pfx) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        
        if level == 0:
            font_size = Pt(18)
            indent = Emu(0)
        elif level == 1:
            font_size = Pt(16)
            indent = Emu(int(0.4 * 914400))
        else:
            font_size = Pt(14)
            indent = Emu(int(0.8 * 914400))
        
        p.level = level
        # Manual indent via left margin
        p.font.size = font_size
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(int(indent)))
        pPr.set('indent', str(int(-0.2 * 914400)) if level > 0 else '0')
        
        if bold_pfx:
            rb = p.add_run()
            rb.text = bold_pfx
            rb.font.name = "Outfit"; rb.font.size = font_size
            rb.font.bold = True; rb.font.color.rgb = INDIGO
        
        r = p.add_run()
        r.text = text
        r.font.name = "Outfit"; r.font.size = font_size
        r.font.bold = False; r.font.color.rgb = GRAY
    
    return slide


def load_caption(fig_dir, fig_name):
    path = os.path.join(fig_dir, f"{fig_name}_caption.txt")
    if os.path.exists(path):
        with open(path) as f:
            return f.read().strip()
    return None


def add_paper_figures(prs, dir_name, figs, titles, authors, journal, year):
    """Add figure slides for a paper. Returns count added."""
    fig_dir = os.path.join(FIG_DIR, dir_name)
    count = 0
    for i, fig_name in enumerate(figs):
        img = os.path.join(fig_dir, f"{fig_name}.png")
        if not os.path.exists(img):
            img = os.path.join(fig_dir, f"{fig_name}.jpg")
        if not os.path.exists(img):
            continue
        title = titles[i] if i < len(titles) else f"Figure {i+1}"
        caption = load_caption(fig_dir, fig_name)
        add_figure_slide(prs, img, title, authors, journal, year, caption)
        count += 1
    return count


# ============================================================
# DECK CONTENT
# ============================================================

def build_deck():
    prs = Presentation(TEMPLATE)
    remove_template_slides(prs)
    n = 0  # slide counter
    
    # === TITLE ===
    add_title_slide(prs, "Neural Alignment",
                    "Comparing DNNs to Brain Activity")
    n += 1
    
    # ============================================================
    # SECTION 1: WHY COMPARE?
    # ============================================================
    add_section_slide(prs, "Why Compare DNNs and Brains?",
                      "Motivation & Framework")
    n += 1
    
    add_text_slide(prs, "The Central Question", [
        (0, "Why would we compare artificial neural networks to biological brains?", ""),
        (0, "", ""),
        (0, "Brain-inspired architectures drove key advances (convolutions, attention, recurrence)", "Engineering: "),
        (0, "DNNs as quantitative, testable models of neural computation", "Neuroscience: "),
        (0, "Understanding the computational principles that give rise to intelligence", "Cognitive science: "),
        (0, "", ""),
        (0, "Neural alignment asks: do DNN internal representations match brain representations?", ""),
    ])
    n += 1
    
    add_text_slide(prs, "Levels of Analysis (Marr, 1982)", [
        (0, "What problem does the system solve? (e.g., object recognition)", "Computational: "),
        (0, "What representations and algorithms are used?", "Algorithmic: "),
        (0, "How is it physically realized? (neurons, silicon)", "Implementational: "),
        (0, "", ""),
        (0, "DNNs primarily operate at the algorithmic level", ""),
        (0, "Neural alignment compares representations across systems that solve the same tasks", ""),
        (0, "Key insight: we can test whether task optimization leads to brain-like solutions", ""),
    ])
    n += 1
    
    # Kriegeskorte & Douglas 2018
    n += add_paper_figures(prs, 'kriegeskorte_2018',
        ['fig_01', 'fig_02', 'fig_03'],
        [f'Cognitive computational neuroscience (Fig. {i})' for i in range(1,4)],
        'Kriegeskorte & Douglas', 'Nature Neurosci.', '2018')
    
    add_text_slide(prs, "The Neuroconnectionist Programme", [
        (0, "Doerig et al. (2023) propose six desiderata for ANN-brain models:", ""),
        (0, "", ""),
        (0, "Explain neurons, areas, and behavior across scales", "1. Multi-scale: "),
        (0, "Match biological wiring constraints", "2. Connectivity: "),
        (0, "Process temporal dynamics like brains", "3. Temporal: "),
        (0, "Use biologically plausible learning rules", "4. Learning: "),
        (0, "Transfer to novel stimuli and tasks", "5. Generalization: "),
        (0, "Capture higher-level cognitive phenomena", "6. Cognitive: "),
    ])
    n += 1
    
    # Doerig et al. 2023
    n += add_paper_figures(prs, 'doerig_2023',
        ['fig_01', 'fig_02', 'fig_03', 'fig_04', 'fig_05'],
        [f'The neuroconnectionist research programme (Fig. {i})' for i in range(1,6)],
        'Doerig et al.', 'Nature Rev. Neurosci.', '2023')
    
    # ============================================================
    # SECTION 2: NEURAL DATA TYPES
    # ============================================================
    add_section_slide(prs, "Neural Data Types & Trade-offs",
                      "fMRI, MEG, ECoG, Single Units")
    n += 1
    
    add_text_slide(prs, "Measuring Brain Activity", [
        (0, "Every neural measurement involves a trade-off between spatial and temporal resolution", ""),
        (0, "", ""),
        (0, "~1mm spatial, ~1s temporal. BOLD signal (blood oxygenation). Whole-brain. Non-invasive.", "fMRI: "),
        (0, "~1cm spatial, ~1ms temporal. Electromagnetic fields. Excellent for dynamics.", "MEG/EEG: "),
        (0, "Single neurons. Invasive. Gold standard for primates (V1, V4, IT).", "Electrophysiology: "),
        (0, "Direct cortical surface. High spatiotemporal. Limited to surgical patients.", "ECoG/iEEG: "),
        (0, "Two-photon. Mouse V1. Thousands of neurons simultaneously.", "Calcium imaging: "),
        (0, "", ""),
        (0, "Each modality constrains what questions we can ask about model-brain alignment", ""),
    ])
    n += 1
    
    # Logothetis 2008
    n += add_paper_figures(prs, 'logothetis_2008',
        ['fig_01', 'fig_02', 'fig_03'],
        [f'What we can do and what we cannot do with fMRI (Fig. {i})' for i in range(1,4)],
        'Logothetis', 'Nature', '2008')
    
    add_text_slide(prs, "Spatiotemporal Fusion", [
        (0, "No single method captures both where and when the brain processes information", ""),
        (0, "", ""),
        (0, "Cichy et al. (2014) pioneered fMRI+MEG fusion through RSA:", ""),
        (1, "Record fMRI and MEG from separate sessions with same stimuli", ""),
        (1, "Build RDMs at each time point (MEG) and each brain region (fMRI)", ""),
        (1, "Cross-correlate to get spatiotemporal dynamics of object recognition", ""),
        (0, "", ""),
        (0, "Result: early visual cortex responds first (~100ms), IT cortex later (~150ms)", "Key finding: "),
    ])
    n += 1
    
    # Cichy et al. 2014
    n += add_paper_figures(prs, 'cichy_2014',
        [f'fig_{i:02d}' for i in range(1,7)],
        [f'Resolving human object recognition in space and time (Fig. {i})' for i in range(1,7)],
        'Cichy, Pantazis & Oliva', 'Nature Neurosci.', '2014')
    
    add_text_slide(prs, "The Natural Scenes Dataset (NSD)", [
        (0, "Allen et al. (2022) created the largest fMRI vision dataset:", ""),
        (0, "", ""),
        (0, "8 subjects, 30,000-40,000 trials each, 7T fMRI", "Scale: "),
        (0, "10,000 natural scene images from COCO", "Stimuli: "),
        (0, "1.8mm isotropic voxels, whole-brain coverage", "Resolution: "),
        (0, "", ""),
        (0, "De facto standard for fMRI-based neural alignment", "Why it matters: "),
        (0, "Built-in noise ceiling estimation (repeat presentations)", ""),
        (0, "Enables fair comparison: same data, same evaluation for all models", ""),
    ])
    n += 1
    
    # Allen et al. 2022
    n += add_paper_figures(prs, 'allen_2022',
        [f'fig_{i:02d}' for i in range(1,7)],
        [f'Natural Scenes Dataset (Fig. {i})' for i in range(1,7)],
        'Allen et al.', 'Nature Neurosci.', '2022')
    
    # ============================================================
    # SECTION 3: METHODS PRIMER
    # ============================================================
    add_section_slide(prs, "Methods Primer",
                      "Encoding, RSA, Brain-Score")
    n += 1
    
    add_text_slide(prs, "Three Approaches to Model-Brain Comparison", [
        (0, "Predict brain activity FROM model features", "1. Encoding models: "),
        (1, "Extract DNN layer activations for each stimulus", ""),
        (1, "Learn linear mapping: Y_brain = W \u00d7 X_model + \u03b5", ""),
        (1, "Ridge regression to prevent overfitting; cross-validate", ""),
        (0, "", ""),
        (0, "Compare representational geometry (no fitting!)", "2. RSA: "),
        (1, "Compute Representational Dissimilarity Matrices (RDMs)", ""),
        (1, "Compare RDMs between model and brain (Spearman correlation)", ""),
        (0, "", ""),
        (0, "Predict stimulus/model FROM brain activity", "3. Decoding: "),
        (1, "Tests what information is present in neural population", ""),
        (1, "Inverse of encoding; different interpretation", ""),
    ])
    n += 1
    
    # Naselaris et al. 2011
    n += add_paper_figures(prs, 'naselaris_2011',
        [f'fig_{i:02d}' for i in range(1,5)],
        [f'Encoding and decoding in fMRI (Fig. {i})' for i in range(1,5)],
        'Naselaris et al.', 'NeuroImage', '2011')
    
    add_text_slide(prs, "RSA: Representational Similarity Analysis", [
        (0, "Kriegeskorte et al. (2008) introduced RSA as a \"common language\" for comparison:", ""),
        (0, "", ""),
        (0, "Present N stimuli \u2192 record brain responses \u2192 compute N\u00d7N distance matrix", "Step 1: "),
        (0, "Present same N stimuli to DNN \u2192 extract layer activations \u2192 compute N\u00d7N distance matrix", "Step 2: "),
        (0, "Correlate the two RDMs (second-order comparison)", "Step 3: "),
        (0, "", ""),
        (0, "Works across species (human fMRI vs. macaque electrophysiology)", "Key advantage: "),
        (0, "Works across measurement types (no need to match dimensions)", ""),
        (0, "Tests representational geometry, not individual unit responses", ""),
    ])
    n += 1
    
    # Kriegeskorte 2008 RSA
    n += add_paper_figures(prs, 'kriegeskorte_2008_rsa',
        [f'fig_{i:02d}' for i in range(1,12)],
        [f'Representational similarity analysis (Fig. {i})' for i in range(1,12)],
        'Kriegeskorte, Mur & Bandettini', 'Front. Syst. Neurosci.', '2008')
    
    add_text_slide(prs, "Noise Ceilings & Benchmarking", [
        (0, "Raw R\u00b2 between model and brain is hard to interpret alone", "Problem: "),
        (0, "A perfect model can't predict beyond the brain's own reliability", ""),
        (0, "", ""),
        (0, "Estimate noise ceiling from data reliability:", "Solution: "),
        (1, "Split-half: correlate odd vs. even trial responses", ""),
        (1, "Spearman-Brown correction: estimate full-data reliability", ""),
        (1, "Report: R\u00b2_normalized = R\u00b2_model / R\u00b2_ceiling", ""),
        (0, "", ""),
        (0, "Brain-Score (Schrimpf et al. 2020) combines neural + behavioral metrics", "Benchmarking: "),
        (1, "Composite score across V1, V2, V4, IT, and behavior", ""),
        (1, "Noise-ceiling normalized \u2192 comparable across datasets", ""),
    ])
    n += 1
    
    # Schrimpf 2020 (Brain-Score)
    n += add_paper_figures(prs, 'schrimpf_2020',
        ['fig_01', 'fig_02'],
        ['Brain-Score (Fig. 1)', 'Brain-Score (Fig. 2)'],
        'Schrimpf et al.', 'Neuron', '2020')
    
    # Nili et al. 2014
    n += add_paper_figures(prs, 'nili_2014',
        ['fig_01', 'fig_02'],
        ['A toolbox for RSA (Fig. 1)', 'A toolbox for RSA (Fig. 2)'],
        'Nili et al.', 'PLOS Comput. Biol.', '2014')
    
    # ============================================================
    # SECTION 4: THE BREAKTHROUGH
    # ============================================================
    add_section_slide(prs, "The Breakthrough",
                      "DNNs Predict Visual Cortex")
    n += 1
    
    add_text_slide(prs, "The Key Insight (2014)", [
        (0, "Hand-designed features (Gabors, SIFT, HMAX) gave moderate neural prediction", "Before: "),
        (0, "Task-optimized DNNs trained on ImageNet spontaneously develop brain-like representations", "After: "),
        (0, "", ""),
        (0, "Yamins et al. (2014) — the paper that launched the field:", ""),
        (1, "Hierarchical modular model trained on object categorization", ""),
        (1, "Predicted V4 and IT neural responses better than any prior model", ""),
        (1, "No neural data used during training \u2014 alignment emerges from the task", ""),
        (0, "", ""),
        (0, "The optimization objective shapes internal representations", "Why it matters: "),
    ])
    n += 1
    
    # Yamins et al. 2014
    n += add_paper_figures(prs, 'yamins_2014',
        [f'fig_{i:02d}' for i in range(1,5)],
        [f'Performance-optimized hierarchical models (Fig. {i})' for i in range(1,5)],
        'Yamins et al.', 'PNAS', '2014')
    
    add_text_slide(prs, "The Complexity Gradient", [
        (0, "DNN layers map systematically onto the visual hierarchy:", ""),
        (0, "", ""),
        (0, "Oriented edges, spatial frequency", "Early layers (conv1-2) \u2192 V1: "),
        (0, "Textures, curvature, moderate complexity", "Middle layers (conv3-4) \u2192 V2/V4: "),
        (0, "Object parts, categories, invariant representations", "Deep layers (fc6-7) \u2192 IT: "),
        (0, "", ""),
        (0, "G\u00fc\u00e7l\u00fc & van Gerven (2015) confirmed this using encoding models on fMRI", ""),
        (0, "Suggests deep networks learn a similar computational progression as the brain", ""),
    ])
    n += 1
    
    # Güçlü & van Gerven 2015
    n += add_paper_figures(prs, 'guclu_2015',
        [f'fig_{i:02d}' for i in range(1,7)],
        [f'Deep neural networks reveal a gradient (Fig. {i})' for i in range(1,7)],
        'G\u00fc\u00e7l\u00fc & van Gerven', 'J. Neurosci.', '2015')
    
    # ============================================================
    # SECTION 5: BEYOND VISION
    # ============================================================
    add_section_slide(prs, "Beyond Vision",
                      "Language & Auditory")
    n += 1
    
    add_text_slide(prs, "Language Model-Brain Alignment", [
        (0, "Same paradigm applied to language:", ""),
        (1, "Present sentences/stories \u2192 record brain activity (fMRI or ECoG)", ""),
        (1, "Extract activations from language models (GPT-2, BERT, etc.)", ""),
        (1, "Map model layers to language brain regions", ""),
        (0, "", ""),
        (0, "Schrimpf et al. (2021) — key findings:", ""),
        (1, "Transformer models predict language cortex best (vs. RNNs, bag-of-words)", ""),
        (1, "Middle-to-late layers align best with temporal/frontal language areas", ""),
        (1, "Next-word prediction as an organizing principle (predictive processing)", ""),
    ])
    n += 1
    
    # Schrimpf 2021 (Language)
    n += add_paper_figures(prs, 'schrimpf_2021_lang',
        [f'fig_{i:02d}' for i in range(1,7)],
        [f'The neural architecture of language (Fig. {i})' for i in range(1,7)],
        'Schrimpf et al.', 'PNAS', '2021')
    
    add_text_slide(prs, "Auditory Neural Alignment", [
        (0, "Kell et al. (2018) extended the Yamins/DiCarlo approach to audition:", ""),
        (0, "", ""),
        (0, "DNN trained on two tasks: word recognition + music genre classification", "Architecture: "),
        (0, "Early shared layers, then two task-specific branches", ""),
        (0, "", ""),
        (0, "DNN representations predicted auditory cortex responses", "Key findings: "),
        (0, "Task-specific branches mapped onto distinct cortical regions", ""),
        (0, "Replicates the vision story: task optimization \u2192 brain-like representations", ""),
    ])
    n += 1
    
    # Kell et al. 2018
    n += add_paper_figures(prs, 'kell_2018',
        [f'fig_{i:02d}' for i in range(1,8)],
        [f'A task-optimized neural network replicates human auditory behavior (Fig. {i})' for i in range(1,8)],
        'Kell et al.', 'Neuron', '2018')
    
    # ============================================================
    # SECTION 6: DECODING
    # ============================================================
    add_section_slide(prs, "Decoding: Brain \u2192 Model",
                      "Reconstruction from Brain Activity")
    n += 1
    
    add_text_slide(prs, "Brain-to-Image Reconstruction", [
        (0, "Takagi & Nishimoto (2023) — fMRI \u2192 images via Stable Diffusion:", ""),
        (0, "", ""),
        (0, "fMRI \u2192 CLIP image embeddings (high-level semantics)", "Semantic pathway: "),
        (0, "fMRI \u2192 VAE latent space (low-level visual features)", "Perceptual pathway: "),
        (0, "Both feed into Stable Diffusion to generate reconstructed image", "Fusion: "),
        (0, "", ""),
        (0, "Visual cortex encodes enough information to reconstruct recognizable scenes", "Implication: "),
        (0, "Not just category \u2014 spatial layout, color, and texture are preserved", ""),
    ])
    n += 1
    
    # Takagi & Nishimoto 2023
    n += add_paper_figures(prs, 'takagi_2023',
        [f'fig_{i:02d}' for i in range(1,4)],
        [f'High-resolution image reconstruction with latent diffusion (Fig. {i})' for i in range(1,4)],
        'Takagi & Nishimoto', 'CVPR', '2023')
    
    add_text_slide(prs, "Brain-to-Text Decoding", [
        (0, "Tang et al. (2023) — continuous language decoding from fMRI:", ""),
        (0, "", ""),
        (0, "Subject listens to stories in scanner", "Setup: "),
        (0, "GPT-based decoder maps fMRI \u2192 word sequences", "Method: "),
        (0, "Uses temporal context (not single-timepoint classification)", ""),
        (0, "", ""),
        (0, "Generates paraphrases (captures meaning, not exact words)", "Result: "),
        (0, "Works across cortex \u2014 not just language areas", ""),
        (0, "Requires subject cooperation (can be resisted)", "Ethics: "),
    ])
    n += 1
    
    # Tang et al. 2023
    n += add_paper_figures(prs, 'tang_2023',
        [f'fig_{i:02d}' for i in range(1,5)],
        [f'Semantic reconstruction of continuous language (Fig. {i})' for i in range(1,5)],
        'Tang et al.', 'Nature Neurosci.', '2023')
    
    # ============================================================
    # SECTION 7: LIMITATIONS
    # ============================================================
    add_section_slide(prs, "Limitations & The Divergence Problem",
                      "What Neural Predictivity Actually Means")
    n += 1
    
    add_text_slide(prs, "The Divergence Problem", [
        (0, "Better ImageNet accuracy \u2192 better brain model", "The assumption: "),
        (0, "", ""),
        (0, "Linsley, Feng & Serre (2025):", "The reality: "),
        (1, "Up to ~70% ImageNet accuracy: performance and brain prediction improve together", ""),
        (1, "Beyond 70%: the correlation breaks down or reverses", ""),
        (1, "State-of-the-art models (ViTs, ConvNeXt) are WORSE brain models than older CNNs", ""),
        (0, "", ""),
        (0, "Models achieve superhuman performance through strategies that diverge from biology", "Why? "),
        (0, "This challenges the entire \"scale up and align\" paradigm", ""),
    ])
    n += 1
    
    # Linsley et al. 2025
    n += add_paper_figures(prs, 'linsley_2025',
        [f'fig_{i:02d}' for i in range(1,4)],
        [f'Better AI does not mean better models of biology (Fig. {i})' for i in range(1,4)],
        'Linsley, Feng & Serre', 'TiCS', '2025')
    
    add_text_slide(prs, "Can Neuroscience Methods Yield Understanding?", [
        (0, "Jonas & Kording (2017) applied standard neuroscience analyses to a microprocessor:", ""),
        (0, "", ""),
        (0, "Lesion studies \u2192 \"transistors critical for Space Invaders\"", ""),
        (0, "Tuning curves \u2192 \"oscillating transistors\"", ""),
        (0, "Dimensionality reduction \u2192 \"low-dimensional manifolds\"", ""),
        (0, "Granger causality \u2192 \"functional connectivity map\"", ""),
        (0, "", ""),
        (0, "All analyses \"work\" but reveal nothing about the actual computation", "Result: "),
        (0, "Correlation and prediction \u2260 understanding", "Lesson: "),
    ])
    n += 1
    
    # Jonas & Kording 2017
    n += add_paper_figures(prs, 'jonas_2017',
        [f'fig_{i:02d}' for i in range(1,14)],
        [f'Could a neuroscientist understand a microprocessor? (Fig. {i})' for i in range(1,14)],
        'Jonas & Kording', 'PLOS Comput. Biol.', '2017')
    
    add_text_slide(prs, "Is Neural Predictivity Diagnostic?", [
        (0, "Storrs et al. (2021) tested whether neural predictivity distinguishes models:", ""),
        (0, "", ""),
        (0, "Trained diverse architectures: AlexNet, VGG, ResNet, SqueezeNet, etc.", ""),
        (0, "After linear reweighting, ALL predict human IT cortex equally well", ""),
        (0, "", ""),
        (0, "If very different models fit the data equally, what does \"brain-like\" mean?", "Challenge: "),
        (0, "Linear fitting may be too flexible \u2014 it can remap any representation", ""),
        (0, "We may need stronger tests than correlation to distinguish models of the brain", ""),
    ])
    n += 1
    
    # Storrs et al. 2021
    n += add_paper_figures(prs, 'storrs_2021',
        [f'fig_{i:02d}' for i in range(1,10)],
        [f'Diverse deep neural networks all predict human IT (Fig. {i})' for i in range(1,10)],
        'Storrs et al.', 'J. Cogn. Neurosci.', '2021')
    
    # CLOSING: Lightning Talk Preview
    add_section_slide(prs, "Thursday: Lightning Talks",
                      "Vision, Language, Audition & Foundation Models")
    n += 1
    
    add_text_slide(prs, "Lightning Talk Preview", [
        (0, "", "Vision:"),
        (1, "Yamins et al. (2014) \u2014 DNNs predict V4/IT", ""),
        (1, "Schrimpf et al. (2020) \u2014 Brain-Score benchmark", ""),
        (1, "Storrs et al. (2021) \u2014 Diverse DNNs all predict IT equally", ""),
        (0, "", "Language:"),
        (1, "Schrimpf et al. (2021) \u2014 Transformers predict language cortex", ""),
        (1, "Caucheteux et al. (2023) \u2014 Predictive coding hierarchy in speech", ""),
        (1, "Goldstein et al. (2022) \u2014 Shared principles: humans & GPT-2", ""),
        (1, "Tuckute et al. (2024) \u2014 Driving/suppressing the language network", ""),
        (0, "", "Audition & Foundation Models:"),
        (1, "Li et al. (2023) \u2014 DNNs for auditory pathway", ""),
        (1, "D\u00e9fossez et al. (2023) \u2014 Decoding speech from MEG", ""),
        (1, "Azabou et al. (2024) \u2014 Unified neural population decoding", ""),
        (1, "Ding et al. (2025) \u2014 Foundation model of mouse visual cortex", ""),
    ])
    n += 1
    
    add_text_slide(prs, "Open Questions", [
        (0, "Is neural predictivity the right metric? (Storrs et al.)", "1. "),
        (0, "What makes a model a good model of the brain vs. a good predictor?", "2. "),
        (0, "Can we go beyond correlation to mechanistic explanation?", "3. "),
        (0, "How do we handle the many-to-one problem? (many models fit equally well)", "4. "),
        (0, "Role of inductive biases: architecture, training data, objectives", "5. "),
        (0, "Does the divergence problem apply beyond vision? (language, audition)", "6. "),
        (0, "", ""),
        (0, "These are the questions your reflection papers should engage with.", ""),
    ])
    n += 1
    
    prs.save(OUTPUT)
    print(f"Saved {n + sum(1 for s in prs.slides)} total slides")
    # Count actual slides
    print(f"Actual slide count: {len(prs.slides)}")


build_deck()
